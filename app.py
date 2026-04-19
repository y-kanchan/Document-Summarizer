"""
DocumentLens NLP Summarizer — Enhanced Engine v2
TextRank + TF-IDF + MMR for local; Claude AI for premium mode.
"""

from flask import Flask, request, jsonify, render_template
import nltk, re, math, os, json, io
from collections import Counter
import pypdf
import docx
import pptx
from werkzeug.utils import secure_filename

# ── NLTK bootstrap ─────────────────────────────────────────────────────────
for pkg in ['punkt', 'punkt_tab', 'stopwords', 'averaged_perceptron_tagger',
            'averaged_perceptron_tagger_eng']:
    try:
        nltk.download(pkg, quiet=True)
    except Exception:
        pass

from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords

app = Flask(__name__)
STOP_WORDS = set(stopwords.words('english'))

# ── File Extraction ──────────────────────────────────────────────────────────

def extract_text_from_file(file_stream, filename):
    ext = filename.rsplit('.', 1)[-1].lower()
    text = ""
    
    if ext == 'pdf':
        reader = pypdf.PdfReader(file_stream)
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text += content + "\n"
    elif ext == 'docx':
        doc = docx.Document(file_stream)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif ext in ['pptx']:
        prs = pptx.Presentation(file_stream)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif ext == 'txt':
        text = file_stream.read().decode('utf-8', errors='ignore')
    
    return text.strip()

# ── Preprocessing ──────────────────────────────────────────────────────────

def preprocess_text(text: str) -> str:
    text = re.sub(r'\r\n|\r', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'([.!?])([A-Z])', r'\1 \2', text)
    text = re.sub(r'[^\x20-\x7E\n]', ' ', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()

def clean_sentence(sent: str) -> str:
    sent = sent.strip()
    sent = re.sub(r'\s+', ' ', sent)
    if sent:
        sent = sent[0].upper() + sent[1:]
    if sent and sent[-1] not in '.!?':
        sent += '.'
    sent = re.sub(r'\s([.,!?;:])', r'\1', sent)
    sent = re.sub(r'([.!?]){2,}', r'\1', sent)
    return sent

def split_sentences(text: str) -> list:
    raw = sent_tokenize(text)
    cleaned = []
    for s in raw:
        s = s.strip()
        if len(s.split()) < 4:
            continue
        if re.match(r'^[\d\W]+$', s):
            continue
        cleaned.append(clean_sentence(s))
    return cleaned

# ── TextRank ───────────────────────────────────────────────────────────────

def sentence_similarity(s1: str, s2: str) -> float:
    def tok(s):
        return set(w.lower() for w in word_tokenize(s)
                   if w.isalpha() and w.lower() not in STOP_WORDS and len(w) > 2)
    w1, w2 = tok(s1), tok(s2)
    if not w1 or not w2:
        return 0.0
    denom = math.log(len(w1) + 1) + math.log(len(w2) + 1)
    return len(w1 & w2) / denom if denom else 0.0

def textrank(sentences: list, damping: float = 0.85, iters: int = 30) -> dict:
    n = len(sentences)
    if n == 0:
        return {}
    sim = [[sentence_similarity(sentences[i], sentences[j]) if i != j else 0.0
            for j in range(n)] for i in range(n)]
    for i in range(n):
        total = sum(sim[i])
        if total > 0:
            sim[i] = [v / total for v in sim[i]]
    scores = [1.0 / n] * n
    for _ in range(iters):
        scores = [(1 - damping) / n + damping * sum(sim[j][i] * scores[j] for j in range(n))
                  for i in range(n)]
    return {sentences[i]: scores[i] for i in range(n)}

# ── TF-IDF ─────────────────────────────────────────────────────────────────

def tfidf_score(sentences: list) -> dict:
    tf = []
    for sent in sentences:
        words = [w.lower() for w in word_tokenize(sent)
                 if w.isalpha() and w.lower() not in STOP_WORDS and len(w) > 2]
        tf.append(Counter(words))
    df = Counter()
    for counter in tf:
        df.update(counter.keys())
    N = len(sentences)
    scores = {}
    for i, sent in enumerate(sentences):
        if not tf[i]:
            scores[sent] = 0.0
            continue
        score = sum((1 + math.log(cnt)) * math.log((N + 1) / (df[w] + 1))
                    for w, cnt in tf[i].items()) / math.sqrt(len(tf[i]))
        scores[sent] = score
    mx = max(scores.values(), default=1) or 1
    return {k: v / mx for k, v in scores.items()}

# ── Position Score ─────────────────────────────────────────────────────────

def position_score(sentences: list) -> dict:
    n = len(sentences)
    scores = {}
    for i, sent in enumerate(sentences):
        p = i / max(n - 1, 1)
        if i == 0:
            scores[sent] = 1.0
        elif i == n - 1:
            scores[sent] = 0.7
        elif p <= 0.15:
            scores[sent] = 0.9 - p * 2
        elif p >= 0.85:
            scores[sent] = 0.5 + (p - 0.85) * 2
        else:
            scores[sent] = 0.2 + 0.1 * math.sin(p * math.pi)
    return scores

# ── MMR Selection ──────────────────────────────────────────────────────────

def mmr_select(sentences: list, base_scores: dict, target: int, lam: float = 0.65) -> list:
    selected = []
    remaining = list(sentences)
    while len(selected) < target and remaining:
        best_sent, best_score = None, -1e9
        for sent in remaining:
            relevance = base_scores.get(sent, 0)
            penalty = max((sentence_similarity(sent, s) for s in selected), default=0.0)
            score = lam * relevance - (1 - lam) * penalty
            if score > best_score:
                best_score, best_sent = score, sent
        if best_sent:
            selected.append(best_sent)
            remaining.remove(best_sent)
    return selected

# ── Local Summarizer ───────────────────────────────────────────────────────

def local_summarize(text: str, ratio: float = 0.30) -> dict:
    text = preprocess_text(text)
    sentences = split_sentences(text)

    if len(sentences) <= 3:
        summary = format_summary(sentences)
        return {
            "summary": summary,
            "original_sentences": len(sentences),
            "summary_sentences": len(sentences),
            "compression_ratio": 1.0,
            "keywords": extract_keywords(text),
            "method": "local-nlp",
        }

    target = max(2, round(len(sentences) * ratio))
    target = min(target, len(sentences))

    tr = textrank(sentences)
    tf = tfidf_score(sentences)
    ps = position_score(sentences)

    base = {s: 0.45 * tr.get(s, 0) + 0.35 * tf.get(s, 0) + 0.20 * ps.get(s, 0)
            for s in sentences}

    selected_set = set(mmr_select(sentences, base, target))
    ordered = [s for s in sentences if s in selected_set]
    summary = format_summary(ordered)

    return {
        "summary": summary,
        "original_sentences": len(sentences),
        "summary_sentences": len(ordered),
        "compression_ratio": round(len(summary) / max(len(text), 1), 3),
        "keywords": extract_keywords(text),
        "method": "local-nlp",
    }

def format_summary(sentences: list) -> str:
    cleaned = [clean_sentence(s) for s in sentences]
    paragraphs = []
    for i in range(0, len(cleaned), 3):
        paragraphs.append(' '.join(cleaned[i:i+3]))
    return '\n\n'.join(paragraphs)

# ── Claude AI Summarizer ───────────────────────────────────────────────────

def claude_summarize(text: str, ratio: float = 0.30, api_key: str = '') -> dict:
    import urllib.request, urllib.error

    length_hint = ("short and concise (1–2 tight paragraphs)" if ratio <= 0.2 else
                   "medium length (3–4 well-developed paragraphs)" if ratio <= 0.4 else
                   "comprehensive (5–6 detailed paragraphs)")

    prompt = (
        "You are a world-class document summarizer. Produce a high-quality, "
        "grammatically flawless, meaningful summary of the document below.\n\n"
        "STRICT RULES:\n"
        "1. Write in clear, flowing prose — NO bullet points, NO headings, NO lists.\n"
        "2. Capture all key ideas, arguments, facts, and conclusions.\n"
        "3. Every sentence must be grammatically correct with perfect punctuation.\n"
        "4. Paragraphs must flow naturally — use transition words.\n"
        "5. Do NOT copy sentences verbatim — synthesise and paraphrase.\n"
        "6. Do NOT add your own opinions or facts not in the document.\n"
        f"7. Target length: {length_hint}.\n"
        "8. Output ONLY the summary — no preamble, no labels.\n\n"
        f"DOCUMENT:\n{text[:14000]}\n\nSUMMARY:"
    )

    payload = json.dumps({
        "model": "claude-sonnet-4-20250514",
        "max_tokens": 1200,
        "messages": [{"role": "user", "content": prompt}]
    }).encode('utf-8')

    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data=payload,
        headers={
            "Content-Type": "application/json",
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01"
        },
        method="POST"
    )

    with urllib.request.urlopen(req, timeout=45) as resp:
        result = json.loads(resp.read().decode('utf-8'))

    summary = result['content'][0]['text'].strip()
    return {
        "summary": summary,
        "original_sentences": len(sent_tokenize(text)),
        "summary_sentences": len(sent_tokenize(summary)),
        "compression_ratio": round(len(summary) / max(len(text), 1), 3),
        "keywords": extract_keywords(text),
        "method": "claude-ai",
    }

# ── Keywords (RAKE-style) ──────────────────────────────────────────────────

def extract_keywords(text: str, top_n: int = 12) -> list:
    stop_pattern = re.compile(
        r'\b(' + '|'.join(re.escape(w) for w in list(STOP_WORDS)[:80]) + r')\b|[.,!?;:()\[\]{}"\']',
        re.IGNORECASE
    )
    phrases = stop_pattern.split(text)
    phrases = [p.strip() for p in phrases if p and len(p.strip()) > 2]

    word_freq: Counter = Counter()
    word_degree: Counter = Counter()
    for phrase in phrases:
        words = [w.lower() for w in word_tokenize(phrase)
                 if w.isalpha() and len(w) > 2 and w.lower() not in STOP_WORDS]
        if not words:
            continue
        deg = len(words) - 1
        for w in words:
            word_freq[w] += 1
            word_degree[w] += deg

    scores = {w: (word_degree[w] + word_freq[w]) / word_freq[w] * math.log(1 + word_freq[w])
              for w in word_freq}
    return sorted(scores, key=scores.get, reverse=True)[:top_n]

# ── Text Analytics ─────────────────────────────────────────────────────────

def _count_syllables(word: str) -> int:
    word = word.lower()
    vowels = "aeiouy"
    count = sum(1 for i, c in enumerate(word)
                if c in vowels and (i == 0 or word[i-1] not in vowels))
    if word.endswith('e') and count > 1:
        count -= 1
    return max(1, count)

def analyze_text(text: str) -> dict:
    sentences = sent_tokenize(text)
    words = [w for w in word_tokenize(text) if w.isalpha()]
    unique = set(w.lower() for w in words)
    avg_len = round(len(words) / max(len(sentences), 1), 1)
    lex_div = round(len(unique) / max(len(words), 1), 3)
    syl = sum(_count_syllables(w) for w in words)
    fre = max(0, min(100, round(
        206.835 - 1.015 * (len(words) / max(len(sentences), 1))
        - 84.6 * (syl / max(len(words), 1)), 1
    ))) if words else 0
    readability = "Easy" if fre >= 70 else "Moderate" if fre >= 50 else "Difficult"
    return {
        "word_count": len(words),
        "sentence_count": len(sentences),
        "unique_words": len(unique),
        "avg_sentence_length": avg_len,
        "lexical_diversity": lex_div,
        "flesch_reading_ease": fre,
        "readability_level": readability,
        "char_count": len(text),
        "read_time_min": max(1, round(len(words) / 200)),
    }

# ── Routes ─────────────────────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/summarize', methods=['POST'])
def summarize():
    import urllib.error
    data = request.get_json()
    if not data or 'text' not in data:
        return jsonify({'error': 'No text provided'}), 400

    text = data['text'].strip()
    if len(text) < 80:
        return jsonify({'error': 'Text too short (min 80 chars).'}), 400
    if len(text) > 60000:
        return jsonify({'error': 'Text too long (max 60,000 chars).'}), 400

    ratio   = max(0.10, min(0.70, float(data.get('ratio', 0.30))))
    api_key = data.get('api_key', '').strip()
    mode    = data.get('mode', 'local')
    warning = None

    try:
        if mode == 'ai' and api_key:
            result = claude_summarize(text, ratio=ratio, api_key=api_key)
        else:
            result = local_summarize(text, ratio=ratio)
    except (urllib.error.URLError, urllib.error.HTTPError, Exception) as e:
        if mode == 'ai':
            result  = local_summarize(text, ratio=ratio)
            warning = f'AI unavailable ({type(e).__name__}), used local NLP.'
        else:
            return jsonify({'error': f'Summarization failed: {str(e)}'}), 500

    stats = analyze_text(text)
    resp = {
        'success': True,
        'summary': result['summary'],
        'stats': stats,
        'summary_stats': {
            'original_sentences': result['original_sentences'],
            'summary_sentences':  result['summary_sentences'],
            'compression_ratio':  result['compression_ratio'],
            'method':             result['method'],
        },
        'keywords': result['keywords'],
    }
    if warning:
        resp['warning'] = warning
    return jsonify(resp)

@app.route('/api/keywords', methods=['POST'])
def keywords():
    data = request.get_json()
    if not data or 'text' not in data:
        return jsonify({'error': 'No text provided'}), 400
    return jsonify({'keywords': extract_keywords(data['text'], int(data.get('top_n', 15)))})

@app.route('/api/analyze', methods=['POST'])
def analyze():
    data = request.get_json()
    if not data or 'text' not in data:
        return jsonify({'error': 'No text provided'}), 400
    return jsonify(analyze_text(data['text']))

@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    ext = file.filename.rsplit('.', 1)[-1].lower()
    if ext not in ['pdf', 'docx', 'pptx', 'txt']:
        return jsonify({'error': f'Unsupported file type: {ext}'}), 400

    try:
        file_stream = io.BytesIO(file.read())
        text = extract_text_from_file(file_stream, file.filename)
        
        if not text:
            return jsonify({'error': 'Could not extract text from file or document is empty.'}), 400
            
        return jsonify({
            'success': True,
            'text': text,
            'filename': file.filename
        })
    except Exception as e:
        return jsonify({'error': f'File processing failed: {str(e)}'}), 500

if __name__ == '__main__':
    print("🚀  DocumentLens v2 — http://127.0.0.1:5000")
    app.run(debug=True, host='0.0.0.0', port=5000)
